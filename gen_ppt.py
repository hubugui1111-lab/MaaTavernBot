"""MaaTavernBot v1.5 —— 自动生成 PPT 文件"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUT_PATH = r"C:\Users\lenovo\Desktop\MaaTavernBot_v15_audio\MaaTavernBot_v1.5.pptx"

# ── 配色（暗色主题，匹配软件界面） ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)       # 深蓝黑背景
BG_CARD = RGBColor(0x25, 0x25, 0x3D)       # 卡片底色
GOLD = RGBColor(0xFF, 0xB3, 0x00)           # 金色标题
WHITE = RGBColor(0xEE, 0xEE, 0xEE)          # 正文白
GRAY = RGBColor(0x99, 0x99, 0xAA)           # 灰色辅助
GREEN = RGBColor(0x4E, 0xC9, 0xB0)          # 绿色高亮
RED = RGBColor(0xFF, 0x6B, 0x6B)            # 红色警告
BLUE = RGBColor(0x5B, 0x9B, 0xF5)           # 蓝色链接

W = Inches(13.333)  # 16:9 宽屏
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ── 辅助函数 ──
def add_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK
    return slide


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_bullets(slide, left, top, width, height, items, size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
    return tf


def add_card(slide, left, top, width, height):
    """添加半透明卡片背景"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════
# 第1页：封面
# ═══════════════════════════════════════════
s = add_slide()
add_text(s, 1.5, 1.5, 10, 1.5, "MaaTavernBot", size=56, color=GOLD, bold=True)
add_text(s, 1.5, 3.0, 10, 0.8, "v1.5  ·  炉石传说酒馆战旗全自动脚本", size=28, color=WHITE)
add_text(s, 1.5, 4.0, 10, 0.6, "基于 MaaFramework  ·  开源免费  ·  仅供学习交流", size=18, color=GRAY)
add_text(s, 1.5, 5.5, 10, 0.5, "github.com/hubugui1111-lab/MaaTavernBot", size=16, color=BLUE)

# ═══════════════════════════════════════════
# 第2页：免责声明
# ═══════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.5, 11, 0.8, "⚠️ 免责声明", size=36, color=RED, bold=True)
add_card(s, 0.8, 1.6, 11.5, 5.0)
warnings = [
    "• 本软件通过图像识别和模拟操作实现游戏自动化",
    "• 可能违反游戏服务条款",
    "• 仅供交流学习 MaaFramework 框架使用",
    "• 使用本软件可能导致游戏账号被警告、限制或永久封禁",
    "• 开发者已尽力降低检测风险，但不对任何账号问题承担责任",
    "• 使用本软件即表示您已充分了解并自愿承担上述风险",
]
add_bullets(s, 1.3, 2.0, 10.5, 4.5, warnings, size=20, color=WHITE, spacing=Pt(16))

# ═══════════════════════════════════════════
# 第3页：v1.5 三大更新
# ═══════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.5, 11, 0.8, "🆕 v1.5 三大更新", size=36, color=GOLD, bold=True)

# 卡片1
add_card(s, 0.8, 1.6, 3.7, 5.0)
add_text(s, 1.1, 1.8, 3.2, 0.5, "🔄 任务自动刷新", size=22, color=GREEN, bold=True)
add_bullets(s, 1.1, 2.6, 3.2, 3.5, [
    "对局胜利后自动检查日常任务",
    "OCR 识别零进度任务",
    "自动点击刷新按钮",
    "每天更换一个未完成任务",
], size=14, color=GRAY)

# 卡片2
add_card(s, 4.8, 1.6, 3.7, 5.0)
add_text(s, 5.1, 1.8, 3.2, 0.5, "🎨 全新顶部导航界面", size=22, color=GREEN, bold=True)
add_bullets(s, 5.1, 2.6, 3.2, 3.5, [
    "顶部导航 + 左侧控制 + 右侧日志",
    "三栏布局，操作更直观",
    "系统托盘最小化",
    "后台静默运行",
], size=14, color=GRAY)

# 卡片3
add_card(s, 8.8, 1.6, 3.7, 5.0)
add_text(s, 9.1, 1.8, 3.2, 0.5, "🔍 本地模板匹配导航", size=22, color=GREEN, bold=True)
add_bullets(s, 9.1, 2.6, 3.2, 3.5, [
    "11 种游戏界面自动识别",
    "不依赖任何外部 API",
    "纯本地运行，速度快",
    "卡住自动返回主菜单",
], size=14, color=GRAY)

# ═══════════════════════════════════════════
# 第4页：任务刷新原理
# ═══════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.5, 11, 0.8, "🔄 任务刷新工作流程", size=36, color=GOLD, bold=True)
add_card(s, 0.8, 1.6, 11.5, 5.3)

steps = [
    ("①", "游戏获胜，检测结算界面"),
    ("②", "点击返回按钮，回到主菜单"),
    ("③", "点击任务图标，进入任务中心"),
    ("④", "点击侧边栏，进入任务界面"),
    ("⑤", "截图 → MinerU OCR 读取进度"),
    ("⑥", "找到进度为 0 的日常任务"),
    ("⑦", "点击刷新箭头，更换任务！"),
    ("⑧", "返回主菜单，准备下一局"),
]
for i, (num, desc) in enumerate(steps):
    y = 2.0 + i * 0.55
    add_text(s, 1.3, y, 0.6, 0.4, num, size=20, color=GREEN, bold=True)
    add_text(s, 2.0, y, 9.5, 0.4, desc, size=18, color=WHITE)

# 警告条
add_text(s, 1.0, 6.5, 11, 0.4, "⚠️ 每个账号每天只能刷新一次    千问多模态 AI 作为 OCR 备用方案",
         size=14, color=RED)

# ═══════════════════════════════════════════
# 第5页：安装步骤
# ═══════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.5, 11, 0.8, "📥 安装只需三步", size=36, color=GOLD, bold=True)

steps_data = [
    ("第一步", "安装 MuMu 模拟器", "分辨率 1600×900 横屏\nDPI 240\n开启 ADB 调试"),
    ("第二步", "安装炉石传说", "登录账号\n确认可正常游玩酒馆战旗"),
    ("第三步", "下载 MaaTavernBot", "GitHub Releases 下载压缩包\n解压到任意目录\n双击 exe 即可运行"),
]

for i, (step, title, detail) in enumerate(steps_data):
    x = 0.8 + i * 4.2
    add_card(s, x, 1.6, 3.8, 5.2)
    add_text(s, x + 0.3, 1.8, 3.2, 0.4, step, size=14, color=GRAY)
    add_text(s, x + 0.3, 2.2, 3.2, 0.5, title, size=24, color=WHITE, bold=True)
    add_text(s, x + 0.3, 3.0, 3.2, 3.5, detail, size=16, color=GRAY)

# ═══════════════════════════════════════════
# 第6页：使用步骤
# ═══════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.5, 11, 0.8, "🚀 使用非常简单", size=36, color=GOLD, bold=True)
add_card(s, 0.8, 1.6, 7.5, 5.3)

usage_steps = [
    "① 打开 MuMu 模拟器，启动炉石传说",
    "② 双击 MaaTavernBot.exe",
    "③ 输入 ADB 端口号（默认 16384 或 16416）",
    "④ 点击「连接」按钮",
    "⑤ 点击「开始」",
]
add_bullets(s, 1.3, 2.0, 6.5, 3.0, usage_steps, size=20, color=WHITE, spacing=Pt(20))

# 定时运行卡片
add_card(s, 8.8, 1.6, 3.7, 3.0)
add_text(s, 9.1, 1.8, 3.2, 0.5, "⏰ 定时运行", size=20, color=GREEN, bold=True)
add_bullets(s, 9.1, 2.5, 3.2, 2.0, [
    "勾选「启用定时运行」",
    "设置开始和结束时间",
    "到时间自动启停",
], size=14, color=GRAY)

# ═══════════════════════════════════════════
# 第7页：常见问题
# ═══════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.5, 11, 0.8, "❓ 常见问题 FAQ", size=36, color=GOLD, bold=True)

faq = [
    ("连接没反应？", "检查 ADB 端口是否正确，模拟器是否开启 ADB 调试"),
    ("不买随从 / 不升级？", "确认模拟器分辨率 1600×900 横屏，DPI 240"),
    ("任务刷新失败？", "每个账号每天只能刷新一次任务"),
    ("卡住不动？", "Bot 有卡住检测，同一画面超过 5 分钟自动重启游戏"),
    ("会不会掉分？", "会。本软件不接入打牌 AI，采用纯盲打策略"),
]

for i, (q, a) in enumerate(faq):
    y = 1.6 + i * 1.1
    add_card(s, 0.8, y, 11.5, 0.9)
    add_text(s, 1.2, y + 0.1, 4.0, 0.4, f"Q: {q}", size=18, color=GREEN, bold=True)
    add_text(s, 1.2, y + 0.45, 10.0, 0.4, f"A: {a}", size=16, color=GRAY)

# ═══════════════════════════════════════════
# 第8页：技术栈
# ═══════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.5, 11, 0.8, "🛠 技术栈 & 开源", size=36, color=GOLD, bold=True)

tech_col1 = [
    ("MaaFramework", "ADB 控制 + 图像识别框架"),
    ("OpenCV", "模板匹配 + 颜色检测"),
    ("MinerU OCR", "本地文字识别，读取任务进度"),
]
tech_col2 = [
    ("千问多模态", "AI 视觉备用方案"),
    ("PySide6", "桌面 GUI 界面"),
    ("edge-tts", "微软晓晓免费配音"),
]

for i, (name, desc) in enumerate(tech_col1 + tech_col2):
    col = 0 if i < 3 else 1
    row = i if i < 3 else i - 3
    x = 0.8 + col * 6.0
    y = 1.6 + row * 1.6
    add_card(s, x, y, 5.5, 1.3)
    add_text(s, x + 0.3, y + 0.2, 5.0, 0.4, name, size=20, color=GREEN, bold=True)
    add_text(s, x + 0.3, y + 0.7, 5.0, 0.4, desc, size=14, color=GRAY)

add_text(s, 0.8, 5.8, 11, 0.5, "🔗 github.com/hubugui1111-lab/MaaTavernBot", size=18, color=BLUE)
add_text(s, 0.8, 6.3, 11, 0.5, "⭐ Release v1.5 已发布，欢迎 Star 支持！", size=20, color=GOLD, bold=True)

# ── 保存 ──
prs.save(OUT_PATH)
print(f"PPT 已生成: {OUT_PATH}")
print(f"共 {len(prs.slides)} 页幻灯片")
print("用 PowerPoint 打开 → 文件 → 导出 → 创建视频 → 配合配音 mp3 即可")
